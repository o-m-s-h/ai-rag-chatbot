from io import BytesIO

from pptx import Presentation

def parse_pptx(file_bytes):

    prs = Presentation(
        BytesIO(file_bytes)
    )

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text